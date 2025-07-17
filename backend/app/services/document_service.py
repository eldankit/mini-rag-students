import PyPDF2
from docx import Document
from pptx import Presentation
import io
from typing import List, Dict, Any
from langchain.text_splitter import RecursiveCharacterTextSplitter
from app.core.config import settings


class DocumentService:
    def __init__(self):
        self.text_splitter = RecursiveCharacterTextSplitter(
            chunk_size=1000,
            chunk_overlap=200,
            length_function=len,
            separators=["\n\n", "\n", " ", ""]
        )
    
    async def extract_text_from_file(self, file_content: bytes, file_type: str) -> str:
        """
        Extract text from different file types
        
        Args:
            file_content: File content as bytes
            file_type: File extension (pdf, docx, pptx, txt)
        
        Returns:
            Extracted text
        """
        try:
            if file_type.lower() == 'pdf':
                return await self._extract_from_pdf(file_content)
            elif file_type.lower() == 'docx':
                return await self._extract_from_docx(file_content)
            elif file_type.lower() == 'pptx':
                return await self._extract_from_pptx(file_content)
            elif file_type.lower() == 'txt':
                return await self._extract_from_txt(file_content)
            else:
                raise ValueError(f"Unsupported file type: {file_type}")
                
        except Exception as e:
            raise Exception(f"Failed to extract text from {file_type} file: {str(e)}")
    
    async def _extract_from_pdf(self, file_content: bytes) -> str:
        """Extract text from PDF file"""
        try:
            pdf_reader = PyPDF2.PdfReader(io.BytesIO(file_content))
            text = ""
            for page in pdf_reader.pages:
                text += page.extract_text() + "\n"
            return text.strip()
        except Exception as e:
            raise Exception(f"PDF extraction failed: {str(e)}")
    
    async def _extract_from_docx(self, file_content: bytes) -> str:
        """Extract text from DOCX file"""
        try:
            doc = Document(io.BytesIO(file_content))
            text = ""
            for paragraph in doc.paragraphs:
                text += paragraph.text + "\n"
            return text.strip()
        except Exception as e:
            raise Exception(f"DOCX extraction failed: {str(e)}")
    
    async def _extract_from_pptx(self, file_content: bytes) -> str:
        """Extract text from PPTX file"""
        try:
            prs = Presentation(io.BytesIO(file_content))
            text = ""
            for slide in prs.slides:
                for shape in slide.shapes:
                    if hasattr(shape, "text"):
                        text += shape.text + "\n"
            return text.strip()
        except Exception as e:
            raise Exception(f"PPTX extraction failed: {str(e)}")
    
    async def _extract_from_txt(self, file_content: bytes) -> str:
        """Extract text from TXT file"""
        try:
            return file_content.decode('utf-8').strip()
        except UnicodeDecodeError:
            # Try with different encoding
            return file_content.decode('latin-1').strip()
    
    async def chunk_text(self, text: str, metadata: Dict[str, Any]) -> List[Dict[str, Any]]:
        """
        Split text into chunks for vector storage
        
        Args:
            text: Text to chunk
            metadata: Metadata to attach to each chunk
        
        Returns:
            List of text chunks with metadata
        """
        try:
            # Split text into chunks
            chunks = self.text_splitter.split_text(text)
            
            # Create documents with metadata
            documents = []
            for i, chunk in enumerate(chunks):
                chunk_metadata = metadata.copy()
                chunk_metadata.update({
                    'chunk_index': i,
                    'total_chunks': len(chunks),
                    'chunk_size': len(chunk)
                })
                
                documents.append({
                    'text': chunk,
                    'metadata': chunk_metadata
                })
            
            return documents
            
        except Exception as e:
            raise Exception(f"Failed to chunk text: {str(e)}")
    
    async def process_file(self, file_content: bytes, file_type: str, metadata: Dict[str, Any]) -> List[Dict[str, Any]]:
        """
        Complete file processing pipeline: extract text and chunk it
        
        Args:
            file_content: File content as bytes
            file_type: File extension
            metadata: File metadata
        
        Returns:
            List of text chunks with metadata
        """
        try:
            # Extract text from file
            text = await self.extract_text_from_file(file_content, file_type)
            
            # Add text length to metadata
            metadata['text_length'] = len(text)
            
            # Chunk the text
            chunks = await self.chunk_text(text, metadata)
            
            return chunks
            
        except Exception as e:
            raise Exception(f"File processing failed: {str(e)}")


# Create a singleton instance
document_service = DocumentService() 